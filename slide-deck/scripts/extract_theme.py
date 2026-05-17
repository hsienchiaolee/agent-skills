#!/usr/bin/env python3
"""Extract a design theme from a reference PowerPoint deck.

Analyzes slides for colors, typography, layout patterns, and component styles,
then outputs a theme.json suitable for build_deck.py.
"""

from __future__ import annotations

import argparse
import json
import sys
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400
EMU_PER_PT = 12700

MONO_FONT_KEYWORDS = ("mono", "courier", "consolas", "fira code", "menlo")

SEMANTIC_COLOR_HINTS = {
    "positive": ("green", "teal"),
    "contrast": ("warm", "orange"),
    "negative": ("red",),
    "neutral": ("dark",),
}


def emu_to_inches(emu: int) -> float:
    return round(emu / EMU_PER_INCH, 3)


def emu_to_pt(emu: int) -> int | None:
    if emu is None:
        return None
    return round(emu / EMU_PER_PT)


def rgb_to_hex(rgb: tuple[int, int, int] | None) -> str | None:
    if rgb is None:
        return None
    return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def safe_rgb(color_obj) -> tuple[int, int, int] | None:
    try:
        rgb = color_obj.rgb
        return (rgb[0], rgb[1], rgb[2])
    except (AttributeError, TypeError, ValueError):
        return None


def brightness(color: tuple[int, int, int]) -> int:
    return sum(color)


# ---------------------------------------------------------------------------
# Shape analysis
# ---------------------------------------------------------------------------

def classify_shape(shape, slide_w: float, slide_h: float) -> str:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "placeholder"

    w = emu_to_inches(shape.width) if shape.width else 0
    h = emu_to_inches(shape.height) if shape.height else 0
    is_auto = shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE

    if h < 0.15 and w > 1.0:
        return "accent_bar"
    if is_auto and w > slide_w * 0.3 and h > slide_h * 0.6:
        return "background_panel"
    if is_auto and w > 2.0 and h > 0.8:
        return "card" if _has_corner_radius(shape) else "rectangle"
    if shape.has_text_frame:
        return "textbox"
    return "shape"


def _has_corner_radius(shape) -> bool:
    try:
        return bool(shape.adjustments and shape.adjustments[0] > 0)
    except (IndexError, TypeError):
        return False


def _get_corner_radius(shape) -> float | None:
    try:
        if shape.adjustments and len(shape.adjustments) > 0:
            return round(float(shape.adjustments[0]), 4)
    except (IndexError, TypeError):
        pass
    return None


def extract_shape_info(shape, slide_w: float, slide_h: float) -> dict:
    role = classify_shape(shape, slide_w, slide_h)
    info = {
        "role": role,
        "left": emu_to_inches(shape.left),
        "top": emu_to_inches(shape.top),
        "width": emu_to_inches(shape.width),
        "height": emu_to_inches(shape.height),
        "fill_color": None,
        "texts": [],
    }

    try:
        if hasattr(shape, "fill") and shape.fill.type is not None:
            rgb = safe_rgb(shape.fill.fore_color)
            if rgb:
                info["fill_color"] = rgb
    except (AttributeError, TypeError):
        pass

    if role == "card":
        radius = _get_corner_radius(shape)
        if radius is not None:
            info["corner_radius"] = radius

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            para_size = emu_to_pt(para.font.size) if para.font.size else None
            para_font = para.font.name
            para_color = safe_rgb(para.font.color) if para.font.color else None
            para_bold = para.font.bold

            for run in para.runs:
                text_info = {"text": run.text}
                size = emu_to_pt(run.font.size) if run.font.size else para_size
                if size:
                    text_info["size"] = size
                font = run.font.name or para_font
                if font:
                    text_info["font"] = font
                rgb = safe_rgb(run.font.color) if run.font.color else para_color
                if rgb:
                    text_info["color"] = rgb
                if run.font.bold or para_bold:
                    text_info["bold"] = True
                info["texts"].append(text_info)

    return info


# ---------------------------------------------------------------------------
# Color clustering
# ---------------------------------------------------------------------------

def cluster_colors(
    color_counts: Counter, threshold: int = 30
) -> list[tuple[tuple[int, int, int], int]]:
    """Group similar colors (within Manhattan distance threshold) and return
    (representative, total_count) sorted by frequency."""
    used: set[tuple[int, int, int]] = set()
    clusters = []

    for color in sorted(color_counts, key=color_counts.get, reverse=True):
        if color in used:
            continue
        cluster = [color]
        used.add(color)
        for other in color_counts:
            if other in used:
                continue
            if sum(abs(a - b) for a, b in zip(color, other)) < threshold:
                cluster.append(other)
                used.add(other)
        clusters.append((color, sum(color_counts[c] for c in cluster)))

    return sorted(clusters, key=lambda x: x[1], reverse=True)


# ---------------------------------------------------------------------------
# Presentation analysis
# ---------------------------------------------------------------------------

def analyze_presentation(pptx_path: str) -> dict:
    prs = Presentation(pptx_path)
    slide_w = emu_to_inches(prs.slide_width)
    slide_h = emu_to_inches(prs.slide_height)

    all_shapes: list[dict] = []
    bg_colors: Counter = Counter()
    fill_colors_by_role: dict[str, Counter] = defaultdict(Counter)
    text_colors: Counter = Counter()
    font_sizes: Counter = Counter()
    font_names: Counter = Counter()
    corner_radii: list[float] = []
    accent_bars: list[dict] = []
    cards: list[dict] = []
    margins_left: list[float] = []

    for slide_idx, slide in enumerate(prs.slides):
        try:
            bg_fill = slide.background.fill
            if bg_fill.type is not None:
                rgb = safe_rgb(bg_fill.fore_color)
                if rgb:
                    bg_colors[rgb] += 1
        except (AttributeError, TypeError):
            pass

        for shape in slide.shapes:
            info = extract_shape_info(shape, slide_w, slide_h)
            info["slide_idx"] = slide_idx
            all_shapes.append(info)

            role = info["role"]
            if info["fill_color"]:
                fill_colors_by_role[role][info["fill_color"]] += 1

            if role == "card":
                cards.append(info)
                if "corner_radius" in info:
                    corner_radii.append(info["corner_radius"])
            elif role == "accent_bar":
                accent_bars.append(info)

            if role in ("textbox", "card") and info["left"] > 0.2:
                margins_left.append(info["left"])

            for t in info["texts"]:
                if t.get("color"):
                    text_colors[t["color"]] += 1
                if t.get("size"):
                    font_sizes[t["size"]] += 1
                if t.get("font"):
                    font_names[t["font"]] += 1

    return {
        "slide_width": slide_w,
        "slide_height": slide_h,
        "num_slides": len(prs.slides),
        "bg_colors": bg_colors,
        "fill_colors_by_role": fill_colors_by_role,
        "text_colors": text_colors,
        "font_sizes": font_sizes,
        "font_names": font_names,
        "corner_radii": corner_radii,
        "accent_bars": accent_bars,
        "cards": cards,
        "margins_left": margins_left,
        "all_shapes": all_shapes,
    }


# ---------------------------------------------------------------------------
# Theme inference
# ---------------------------------------------------------------------------

def _mean(values: list[float]) -> float:
    return sum(values) / len(values) if values else 0


def _pick_size_at(sorted_sizes: list[int], fraction: float, default: int) -> int:
    if not sorted_sizes or len(sorted_sizes) <= 1:
        return default
    idx = min(int(len(sorted_sizes) * fraction), len(sorted_sizes) - 1)
    return sorted_sizes[idx]


def infer_color_roles(analysis: dict) -> dict:
    fills = analysis["fill_colors_by_role"]
    colors: dict[str, tuple | list] = {}

    bg_clustered = cluster_colors(analysis["bg_colors"])
    if bg_clustered:
        colors["bg_light"] = max(bg_clustered[:3], key=lambda c: brightness(c[0]))[0]
        dark_candidates = [c for c in bg_clustered if brightness(c[0]) < 400]
        if dark_candidates:
            colors["bg_dark"] = dark_candidates[0][0]

    card_fills = fills.get("card", Counter())
    if card_fills:
        card_clustered = cluster_colors(card_fills)
        if card_clustered:
            colors["card_bg"] = card_clustered[0][0]

    accent_fills = fills.get("accent_bar", Counter())
    if accent_fills:
        colors["accents"] = [
            color
            for color, _ in cluster_colors(accent_fills)[:6]
            if 100 < brightness(color) < 600
        ]

    text_clustered = cluster_colors(analysis["text_colors"])
    heading_candidates = []
    body_candidates = []
    muted_candidates = []
    for color, count in text_clustered:
        b = brightness(color)
        if b < 300:
            heading_candidates.append(color)
        elif b < 500:
            body_candidates.append(color)
        elif b < 650:
            muted_candidates.append(color)
    if heading_candidates:
        colors["text_heading"] = min(heading_candidates, key=brightness)
    if body_candidates:
        colors["text_body"] = body_candidates[0]
    if muted_candidates:
        colors["text_muted"] = muted_candidates[0]

    code_fills = fills.get("rectangle", Counter()) or fills.get("card", Counter())
    for color, _ in cluster_colors(code_fills):
        if brightness(color) < 300:
            colors["code_bg"] = color
            break

    return colors


def name_accents(accent_colors: list[tuple[int, int, int]]) -> dict[str, tuple[int, int, int]]:
    """Assign human-readable names to accent colors based on hue."""
    names: dict[str, tuple[int, int, int]] = {}
    for color in accent_colors:
        r, g, b = color

        if brightness(color) < 200:
            name = "dark"
        elif max(r, g, b) == min(r, g, b):
            name = "gray"
        elif r >= g and r >= b:
            name = "warm" if g > b + 30 else "red"
        elif g >= r and g >= b:
            if b > r:
                name = "teal" if b > g * 0.5 else "green"
            else:
                name = "olive" if r > g * 0.5 else "green"
        else:
            name = "blue" if r < g else "purple"

        base = name
        suffix = 2
        while name in names:
            name = f"{base}{suffix}"
            suffix += 1
        names[name] = color

    return names


def _assign_semantic(accent_names: list[str]) -> dict[str, str]:
    """Map semantic roles (positive, contrast, etc.) to accent names."""
    semantic: dict[str, str] = {}
    for fallback_idx, (role, hints) in enumerate(SEMANTIC_COLOR_HINTS.items()):
        for name in accent_names:
            if any(h in name for h in hints):
                semantic[role] = name
                break
        if role not in semantic and fallback_idx < len(accent_names):
            semantic.setdefault(role, accent_names[fallback_idx])
    return semantic


def infer_typography(analysis: dict) -> dict:
    sizes = analysis["font_sizes"]
    fonts = analysis["font_names"]

    heading_font = fonts.most_common(1)[0][0] if fonts else "Arial"
    mono_candidates = [
        f for f in fonts
        if any(kw in f.lower() for kw in MONO_FONT_KEYWORDS)
    ]
    mono_font = mono_candidates[0] if mono_candidates else "Courier New"

    sorted_sizes = sorted(sizes.keys())
    scale = {}
    if sorted_sizes:
        scale = {
            "badge": min(sorted_sizes),
            "code": _pick_size_at(sorted_sizes, 0.2, 11),
            "small": _pick_size_at(sorted_sizes, 0.25, 11),
            "caption": _pick_size_at(sorted_sizes, 0.25, 11),
            "body": _pick_size_at(sorted_sizes, 0.5, 13),
            "subtitle": _pick_size_at(sorted_sizes, 0.6, 15),
            "card_title": _pick_size_at(sorted_sizes, 0.7, 16),
            "slide_title": _pick_size_at(sorted_sizes, 0.85, 24),
            "hero_title": max(sorted_sizes),
        }

    return {
        "heading_font": heading_font,
        "body_font": heading_font,
        "mono_font": mono_font,
        "scale": scale,
    }


def infer_layout(analysis: dict) -> dict:
    margins = analysis["margins_left"]
    cards = analysis["cards"]

    left_margin = round(min(margins), 1) if margins else 0.8

    card_gaps: list[float] = []
    by_slide: dict[int, list[dict]] = defaultdict(list)
    for c in cards:
        by_slide[c["slide_idx"]].append(c)
    for slide_cards in by_slide.values():
        sorted_cards = sorted(slide_cards, key=lambda c: c["left"])
        for i in range(len(sorted_cards) - 1):
            gap = sorted_cards[i + 1]["left"] - (sorted_cards[i]["left"] + sorted_cards[i]["width"])
            if 0.1 < gap < 2.0:
                card_gaps.append(gap)

    content_tops = [c["top"] for c in cards if c["top"] > 0.5]

    return {
        "margin": {"left": left_margin, "right": left_margin, "top": 0.4, "bottom": 0.4},
        "content_top": round(min(content_tops), 1) if content_tops else 1.3,
        "card_gap": round(_mean(card_gaps), 1) if card_gaps else 0.3,
        "card_padding": 0.4,
        "badge_position": [left_margin, 0.15],
        "subtitle_top": 1.2,
    }


def infer_components(analysis: dict) -> dict:
    radii = analysis["corner_radii"]
    bars = analysis["accent_bars"]

    return {
        "card": {
            "corner_radius": round(_mean(radii), 4) if radii else 0.02,
            "accent_bar": len(bars) > 0,
            "accent_bar_height": round(_mean([b["height"] for b in bars]), 3) if bars else 0.06,
            "border": False,
            "shadow": False,
        },
        "bullet": {
            "character": "•",
            "spacing_pt": 10,
            "sub_character": "›",
            "sub_indent_spaces": 4,
        },
        "callout": {
            "height": 0.6,
            "accent": "dark",
            "bold": True,
            "font_size": 12,
            "bottom_margin": 0.4,
        },
        "badge": {
            "uppercase": True,
            "bold": True,
            "color": "green",
        },
        "code_block": {
            "padding": 0.3,
            "line_spacing_pt": 2,
        },
    }


def infer_slide_patterns(analysis: dict) -> dict:
    shapes = analysis["all_shapes"]
    slide_w = analysis["slide_width"]

    title_pattern: dict = {"split_position": 6.8, "left_bg": "light", "right_bg": "dark"}
    for panel in shapes:
        if panel["role"] == "background_panel" and panel["left"] > slide_w * 0.3:
            title_pattern["split_position"] = round(panel["left"], 1)
            break

    for shape in shapes:
        if shape["role"] == "image":
            title_pattern["image_area"] = {
                "left": shape["left"],
                "top": shape["top"],
                "size": round(max(shape["width"], shape["height"]), 1),
            }
            break

    max_text_size = max(
        (t.get("size", 0) for s in shapes for t in s.get("texts", [])),
        default=0,
    )

    return {
        "title": title_pattern,
        "closing": {
            "background": "dark",
            "title_size": max_text_size if max_text_size > 30 else 56,
            "centered": True,
        },
    }


# ---------------------------------------------------------------------------
# Theme assembly
# ---------------------------------------------------------------------------

def build_theme(analysis: dict) -> dict:
    color_roles = infer_color_roles(analysis)
    typography = infer_typography(analysis)
    layout = infer_layout(analysis)
    components = infer_components(analysis)
    patterns = infer_slide_patterns(analysis)

    accent_colors = color_roles.get("accents", [])
    named_accents = name_accents(accent_colors)
    accent_names = list(named_accents.keys())

    semantic = _assign_semantic(accent_names) if accent_names else {}
    components["badge"]["color"] = semantic.get("positive", accent_names[0] if accent_names else "green")
    components["callout"]["accent"] = semantic.get("neutral", "dark")

    defaults = {
        "bg_light": (255, 255, 255),
        "bg_dark": (27, 58, 75),
        "text_heading": (45, 52, 54),
        "text_body": (95, 107, 109),
        "text_muted": (139, 149, 152),
        "card_bg": (238, 244, 242),
        "code_bg": (27, 58, 75),
    }

    def color(key):
        return rgb_to_hex(color_roles.get(key, defaults[key]))

    return {
        "meta": {
            "name": Path(analysis.get("source", "extracted")).stem + " theme",
            "source": analysis.get("source", ""),
            "slide_width": analysis["slide_width"],
            "slide_height": analysis["slide_height"],
        },
        "colors": {
            "background": {"light": color("bg_light"), "dark": color("bg_dark")},
            "text": {
                "heading": color("text_heading"),
                "body": color("text_body"),
                "muted": color("text_muted"),
                "on_dark": "#FFFFFF",
            },
            "accents": {n: rgb_to_hex(c) for n, c in named_accents.items()},
            "card": {
                "background": color("card_bg"),
                "code_background": color("code_bg"),
                "code_text": "#CDD6F4",
            },
            "semantic": semantic,
        },
        "typography": typography,
        "layout": layout,
        "components": components,
        "slide_patterns": patterns,
    }


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Extract a design theme from a PPTX")
    parser.add_argument("pptx", help="Path to the reference PowerPoint file")
    parser.add_argument("-o", "--output", default="theme.json", help="Output theme file")
    parser.add_argument("--verbose", action="store_true", help="Print analysis details")
    args = parser.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"Error: {pptx_path} not found", file=sys.stderr)
        sys.exit(1)

    print(f"Analyzing {pptx_path}...")
    analysis = analyze_presentation(str(pptx_path))
    analysis["source"] = str(pptx_path)

    if args.verbose:
        print(f"  Slides: {analysis['num_slides']}")
        print(f"  Dimensions: {analysis['slide_width']}\" x {analysis['slide_height']}\"")
        print(f"  Background colors: {len(analysis['bg_colors'])}")
        print(f"  Cards found: {len(analysis['cards'])}")
        print(f"  Accent bars found: {len(analysis['accent_bars'])}")
        print(f"  Font sizes: {sorted(analysis['font_sizes'].keys())}")
        print(f"  Fonts: {list(analysis['font_names'].keys())}")

    theme = build_theme(analysis)

    output_path = Path(args.output)
    with open(output_path, "w") as f:
        json.dump(theme, f, indent=2)

    print(f"Theme saved to {output_path}")
    print(f"  Colors: {len(theme['colors']['accents'])} accents")
    print(f"  Fonts: {theme['typography']['heading_font']}, {theme['typography']['mono_font']}")
    print(f"  Slide size: {theme['meta']['slide_width']}\" x {theme['meta']['slide_height']}\"")


if __name__ == "__main__":
    main()

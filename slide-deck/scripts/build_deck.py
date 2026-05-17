#!/usr/bin/env python3
"""Build a PowerPoint deck from structured markdown and a theme.json.

Usage:
    python build_deck.py slides.md --theme theme.json -o output.pptx
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _is_dark_color(color: RGBColor) -> bool:
    return (color[0] * 299 + color[1] * 587 + color[2] * 114) / 1000 < 128


def _colors_too_similar(a: RGBColor, b: RGBColor, threshold: int = 80) -> bool:
    return abs(a[0] - b[0]) + abs(a[1] - b[1]) + abs(a[2] - b[2]) < threshold


_INLINE_RE = re.compile(r"(\*\*(.+?)\*\*|`(.+?)`)")


def _estimate_text_height(text: str, font_size_pt: int, width_inches: float,
                          bold: bool = False) -> float:
    avg_char_width_pt = font_size_pt * (0.65 if bold else 0.55)
    width_pt = width_inches * 72
    chars_per_line = max(1, int(width_pt / avg_char_width_pt))
    num_lines = max(1, -(-len(text) // chars_per_line))
    line_height_pt = font_size_pt * 1.2
    return (num_lines * line_height_pt) / 72


def _parse_inline_markdown(text: str) -> list[dict]:
    parts = []
    last = 0
    for m in _INLINE_RE.finditer(text):
        if m.start() > last:
            parts.append({"text": text[last:m.start()], "style": "normal"})
        if m.group(2):
            parts.append({"text": m.group(2), "style": "bold"})
        elif m.group(3):
            parts.append({"text": m.group(3), "style": "code"})
        last = m.end()
    if last < len(text):
        parts.append({"text": text[last:], "style": "normal"})
    return parts if parts else [{"text": text, "style": "normal"}]


# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------

class Theme:
    def __init__(self, data: dict):
        self._data = data
        self.slide_width: float = data["meta"]["slide_width"]
        self.slide_height: float = data["meta"]["slide_height"]

        colors = data["colors"]
        self.bg_light = hex_to_rgb(colors["background"]["light"])
        self.bg_dark = hex_to_rgb(colors["background"]["dark"])
        self.text_heading = hex_to_rgb(colors["text"]["heading"])
        self.text_body = hex_to_rgb(colors["text"]["body"])
        self.text_muted = hex_to_rgb(colors["text"]["muted"])
        self.text_on_dark = hex_to_rgb(colors["text"]["on_dark"])
        self.card_bg = hex_to_rgb(colors["card"]["background"])
        self.code_bg = hex_to_rgb(colors["card"]["code_background"])
        self.code_text = hex_to_rgb(colors["card"]["code_text"])

        self.accents = {k: hex_to_rgb(v) for k, v in colors.get("accents", {}).items()}
        self.semantic: dict[str, str] = colors.get("semantic", {})

        typo = data["typography"]
        self.heading_font: str = typo["heading_font"]
        self.body_font: str = typo["body_font"]
        self.mono_font: str = typo["mono_font"]
        self.scale: dict[str, int] = typo["scale"]

        layout = data["layout"]
        self.margin_left: float = layout["margin"]["left"]
        self.margin_right: float = layout["margin"]["right"]
        self.margin_top: float = layout["margin"]["top"]
        self.content_top: float = layout["content_top"]
        self.card_gap: float = min(layout["card_gap"], 0.3)
        self.card_padding: float = layout["card_padding"]
        self.badge_pos: list[float] = layout.get("badge_position", [self.margin_left, 0.15])
        self.subtitle_top: float = layout.get("subtitle_top", 1.2)

        self.comp: dict = data["components"]
        self.patterns: dict = data.get("slide_patterns", {})

    def accent(self, name: str | None) -> RGBColor:
        if name and name in self.accents:
            return self.accents[name]
        resolved = self.semantic.get(name, "") if name else ""
        if resolved in self.accents:
            return self.accents[resolved]
        if self.accents:
            return next(iter(self.accents.values()))
        return self.text_heading

    @property
    def content_width(self) -> float:
        return self.slide_width - self.margin_left - self.margin_right

    def font_size(self, key: str, default: int = 13) -> int:
        return self.scale.get(key, default)

    @classmethod
    def load(cls, path: str) -> Theme:
        with open(path) as f:
            return cls(json.load(f))


# ---------------------------------------------------------------------------
# Markdown parser
# ---------------------------------------------------------------------------

def parse_frontmatter(text: str) -> tuple[dict[str, str], str]:
    text = text.strip()
    if not text.startswith("---"):
        return {}, text
    try:
        end = text.index("---", 3)
    except ValueError:
        return {}, text
    fm_text = text[3:end].strip()
    rest = text[end + 3:].strip()
    fm = {}
    for line in fm_text.split("\n"):
        if ":" in line:
            k, v = line.split(":", 1)
            fm[k.strip()] = v.strip()
    return fm, rest


def _parse_heading(line: str) -> tuple[str, list[str]]:
    """Extract heading text and bracket modifiers from a # or ## line."""
    text = line.lstrip("#").strip()
    mods = []
    m = re.search(r"\[([^\]]+)\]\s*$", text)
    if m:
        mods = [x.strip().lower() for x in m.group(1).split(",")]
        text = text[:m.start()].strip()
    return text, mods


def _parse_bullet(line: str) -> dict:
    indent = len(line) - len(line.lstrip())
    text = line.strip().removeprefix("- ").strip()
    return {"type": "bullet", "text": text, "level": 1 if indent >= 2 else 0}


def _parse_metadata(line: str) -> tuple[str, str] | None:
    """Check if line is a metadata field (author:, date:, image:)."""
    for key in ("author", "date", "image"):
        if line.startswith(f"{key}:"):
            return key, line.split(":", 1)[1].strip()
    return None


def parse_slides(text: str) -> list[dict]:
    slides = []
    for raw in re.split(r"\n---\n", text):
        raw = raw.strip()
        if raw:
            slide = _parse_single_slide(raw)
            if slide:
                slides.append(slide)
    return slides


def _parse_single_slide(text: str) -> dict:
    slide = {
        "title": "",
        "mods": [],
        "badge": None,
        "subtitle": None,
        "sections": [],
        "takeaway": None,
        "metadata": {},
        "body_lines": [],
    }

    current_section: dict | None = None
    in_code = False
    code_lines: list[str] = []

    for line in text.split("\n"):
        # Code fence toggle
        if line.startswith("```"):
            if in_code:
                code_item = {"type": "code", "text": "\n".join(code_lines)}
                if current_section:
                    current_section["content"].append(code_item)
                else:
                    slide["sections"].append({
                        "heading": "",
                        "accent": None,
                        "mods": ["code_only"],
                        "content": [code_item],
                    })
                code_lines = []
            in_code = not in_code
            continue

        if in_code:
            code_lines.append(line)
            continue

        # Slide title
        if line.startswith("# ") and not line.startswith("## "):
            slide["title"], slide["mods"] = _parse_heading(line)
            continue

        # Card section
        if line.startswith("## "):
            heading, mods = _parse_heading(line)
            accent = next((m for m in mods if m != "full"), None)
            section_mods = [m for m in mods if m == "full"]
            current_section = {
                "heading": heading,
                "accent": accent,
                "mods": section_mods,
                "content": [],
            }
            slide["sections"].append(current_section)
            continue

        if line.startswith("badge:"):
            slide["badge"] = line.split(":", 1)[1].strip()
            continue

        meta = _parse_metadata(line)
        if meta:
            slide["metadata"][meta[0]] = meta[1]
            continue

        if line.startswith("!>"):
            slide["takeaway"] = line[2:].strip()
            continue

        if line.startswith("> "):
            value = line[2:].strip()
            slide["subtitle"] = f"{slide['subtitle']} {value}" if slide["subtitle"] else value
            continue

        stripped = line.strip()

        if stripped.startswith("- "):
            item = _parse_bullet(line)
            target = current_section["content"] if current_section else slide["body_lines"]
            target.append(item)
            continue

        if stripped.startswith("**") and stripped.endswith("**"):
            if current_section:
                current_section["content"].append({"type": "bold", "text": stripped[2:-2]})
            continue

        if stripped and current_section:
            current_section["content"].append({"type": "text", "text": stripped})

    return slide


# ---------------------------------------------------------------------------
# Shape primitives
# ---------------------------------------------------------------------------

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, color: RGBColor,
               shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_card(slide, left: float, top: float, width: float, height: float,
              theme: Theme, accent_name: str | None = None, *, dark_card: bool = False):
    shape = _add_shape(
        slide, Inches(left), Inches(top), Inches(width), Inches(height),
        theme.card_bg,
    )

    if accent_name and theme.comp["card"].get("accent_bar", True):
        bar_h = theme.comp["card"].get("accent_bar_height", 0.06)
        bar_color = theme.accent(accent_name)
        if dark_card and _colors_too_similar(bar_color, theme.card_bg):
            bar_color = theme.text_on_dark
        _add_shape(slide, Inches(left), Inches(top), Inches(width), Inches(bar_h),
                   bar_color)
    return shape


def _add_textbox(slide, left: float, top: float, width: float, height: float,
                 text: str, font_size: int, color: RGBColor, theme: Theme, *,
                 bold: bool = False, alignment=PP_ALIGN.LEFT, font_name: str | None = None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name or theme.heading_font
    p.alignment = alignment
    return tb


def _add_bullet_list(slide, left: float, top: float, width: float, height: float,
                     items: list[dict], theme: Theme, *,
                     font_size: int | None = None, text_color: RGBColor | None = None,
                     bullet_color: RGBColor | None = None):
    font_size = font_size or theme.font_size("body")
    text_color = text_color or theme.text_body
    bullet_color = bullet_color or theme.accent("positive")
    spacing = Pt(theme.comp["bullet"].get("spacing_pt", 10))
    bullet_char = theme.comp["bullet"].get("character", "•")
    sub_char = theme.comp["bullet"].get("sub_character", "›")
    indent_spaces = theme.comp["bullet"].get("sub_indent_spaces", 4)

    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing

        level = item.get("level", 0)
        char = sub_char if level else bullet_char

        run_b = p.add_run()
        run_b.text = " " * (level * indent_spaces) + char + "  "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = theme.body_font

        item_size = Pt(font_size if level == 0 else font_size - 1)
        for part in _parse_inline_markdown(item["text"]):
            run_t = p.add_run()
            run_t.text = part["text"]
            run_t.font.size = item_size
            run_t.font.color.rgb = text_color
            run_t.font.name = theme.body_font
            if part["style"] == "bold":
                run_t.font.bold = True
            elif part["style"] == "code":
                run_t.font.name = theme.mono_font

    return tb


def _add_code_block(slide, left: float, top: float, width: float, height: float,
                    text: str, theme: Theme):
    corner_radius = theme.comp["card"].get("corner_radius", 0.02)
    shape = _add_shape(
        slide, Inches(left), Inches(top), Inches(width), Inches(height),
        theme.code_bg, MSO_SHAPE.ROUNDED_RECTANGLE,
    )
    shape.adjustments[0] = corner_radius

    pad = theme.comp["code_block"].get("padding", 0.3)
    line_sp = Pt(theme.comp["code_block"].get("line_spacing_pt", 2))
    code_size = theme.font_size("code", 11)

    tb = slide.shapes.add_textbox(
        Inches(left + pad), Inches(top + pad * 0.8),
        Inches(width - pad * 2), Inches(height - pad * 1.6))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = line_sp
        run = p.add_run()
        run.text = line
        run.font.size = Pt(code_size)
        run.font.color.rgb = theme.code_text
        run.font.name = theme.mono_font

    return tb


# ---------------------------------------------------------------------------
# Layout engine
# ---------------------------------------------------------------------------

def _build_title_slide(prs, slide_data: dict, theme: Theme,
                       frontmatter: dict, md_dir: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.bg_light)

    tp = theme.patterns.get("title", {})
    split = tp.get("split_position", 6.8)
    _add_shape(slide, Inches(split), Inches(0),
               Inches(theme.slide_width - split), Inches(theme.slide_height),
               theme.bg_dark)

    ml = theme.margin_left
    left_panel_w = split - ml - 0.5

    badge_text = slide_data.get("badge")
    if badge_text:
        _add_textbox(slide, ml, 1.2, left_panel_w, 0.3,
                     badge_text.upper(), theme.font_size("badge", 9),
                     theme.accent("positive"), theme, bold=True)

    _add_textbox(slide, ml, 1.6, left_panel_w, 1.8,
                 slide_data["title"], theme.font_size("hero_title", 44),
                 theme.text_heading, theme, bold=True)

    if slide_data.get("subtitle"):
        _add_textbox(slide, ml, 3.5, left_panel_w, 0.8,
                     slide_data["subtitle"], theme.font_size("subtitle", 15),
                     theme.text_body, theme)

    author = slide_data["metadata"].get("author") or frontmatter.get("author", "")
    if author:
        _add_textbox(slide, ml, 4.6, left_panel_w, 0.5,
                     author, theme.font_size("subtitle", 15),
                     theme.text_heading, theme, bold=True)

    date = slide_data["metadata"].get("date") or frontmatter.get("date", "")
    if date:
        _add_textbox(slide, ml, 5.1, left_panel_w, 0.4,
                     str(date), theme.font_size("caption", 11),
                     theme.text_muted, theme)

    image_path = slide_data["metadata"].get("image") or frontmatter.get("title_image")
    if image_path:
        img_path = Path(md_dir) / image_path
        if img_path.exists():
            ia = tp.get("image_area", {"size": 1.0})
            img_size = ia.get("size", 1.0)
            right_panel_w = theme.slide_width - split
            img_left = split + (right_panel_w - img_size) / 2
            img_top = (theme.slide_height - img_size) / 2
            slide.shapes.add_picture(
                str(img_path),
                Inches(img_left), Inches(img_top),
                Inches(img_size), Inches(img_size))


def _build_dark_slide(prs, slide_data: dict, theme: Theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.bg_dark)

    cp = theme.patterns.get("closing", {})
    _add_textbox(slide, 1.5, 2.8, theme.slide_width - 3.0, 1.2,
                 slide_data["title"], cp.get("title_size", 56),
                 theme.text_on_dark, theme, bold=True,
                 alignment=PP_ALIGN.CENTER)

    if slide_data.get("subtitle"):
        _add_textbox(slide, 2.0, 4.3, theme.slide_width - 4.0, 0.5,
                     slide_data["subtitle"], theme.font_size("card_title", 18),
                     theme.text_muted, theme, alignment=PP_ALIGN.CENTER)


def _build_content_slide(prs, slide_data: dict, theme: Theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.bg_light)
    ml = theme.margin_left

    title_top = theme.margin_top
    if slide_data.get("badge"):
        badge_comp = theme.comp.get("badge", {})
        badge_top = theme.badge_pos[1] if len(theme.badge_pos) > 1 else 0.15
        _add_textbox(
            slide, ml, badge_top, 2.0, 0.3,
            slide_data["badge"].upper() if badge_comp.get("uppercase", True) else slide_data["badge"],
            badge_comp.get("font_size", theme.font_size("badge", 9)),
            theme.accent(badge_comp.get("color", "green")), theme,
            bold=badge_comp.get("bold", True),
        )
        title_top = max(title_top, 0.5)

    title_size = min(theme.font_size("slide_title", 24), 26)
    title_h = _estimate_text_height(slide_data["title"], title_size, theme.content_width, bold=True)
    _add_textbox(slide, ml, title_top, theme.content_width, title_h,
                 slide_data["title"], title_size,
                 theme.text_heading, theme, bold=True)

    next_top = title_top + title_h + 0.05
    if slide_data.get("subtitle"):
        sub_size = min(theme.font_size("subtitle", 15), 16)
        sub_h = _estimate_text_height(slide_data["subtitle"], sub_size, theme.content_width)
        _add_textbox(slide, ml, next_top, theme.content_width, sub_h,
                     slide_data["subtitle"], sub_size,
                     theme.text_body, theme)
        next_top += sub_h + 0.1

    regular, full, code_only = [], [], []
    for s in slide_data.get("sections", []):
        mods = s.get("mods", [])
        if "code_only" in mods:
            code_only.append(s)
        elif "full" in mods:
            full.append(s)
        else:
            regular.append(s)

    has_code = bool(code_only) or any(
        any(c["type"] == "code" for c in s.get("content", []))
        for s in regular
    )

    card_top = max(next_top + 0.1, theme.content_top)
    callout_cfg = theme.comp.get("callout", {})
    takeaway_h = callout_cfg.get("height", 0.6) if slide_data.get("takeaway") else 0
    bottom_margin = callout_cfg.get("bottom_margin", 0.4)
    gap = theme.card_gap
    takeaway_gap = gap if slide_data.get("takeaway") else 0
    available_h = theme.slide_height - card_top - takeaway_h - bottom_margin - takeaway_gap

    if full:
        row_gap = gap * 0.5
        top_h = available_h * 0.38
        bottom_h = available_h - top_h - row_gap
        if regular:
            _place_cards(slide, regular, code_only, has_code, card_top, top_h, theme)
        bottom_top = card_top + top_h + row_gap
        for i, fs in enumerate(full):
            fw = (theme.content_width - gap * (len(full) - 1)) / len(full)
            fl = ml + i * (fw + gap)
            _render_card(slide, fs, fl, bottom_top, fw, bottom_h, theme)
    else:
        _place_cards(slide, regular, code_only, has_code, card_top, available_h, theme)

    if slide_data.get("body_lines") and not regular and not full:
        bullets = [b for b in slide_data["body_lines"] if b.get("type") == "bullet"]
        if bullets:
            _add_bullet_list(slide, ml, card_top, theme.content_width, available_h,
                             bullets, theme)

    if slide_data.get("takeaway"):
        _render_takeaway(slide, slide_data["takeaway"], theme)



def _place_cards(slide, regular: list[dict], code_only: list[dict],
                 has_code: bool, top: float, height: float, theme: Theme):
    """Dispatch to the right column layout based on section count and content."""
    if has_code and len(regular) <= 1:
        _layout_card_and_code(slide, regular, code_only, top, height, theme)
    elif len(regular) == 1:
        _render_card(slide, regular[0], theme.margin_left, top, theme.content_width, height, theme)
    elif len(regular) >= 2:
        _layout_columns(slide, regular, top, height, theme)


def _layout_columns(slide, sections: list[dict], top: float, height: float, theme: Theme):
    ml = theme.margin_left
    gap = theme.card_gap
    n = len(sections)
    card_w = (theme.content_width - gap * (n - 1)) / n
    for i, sec in enumerate(sections):
        _render_card(slide, sec, ml + i * (card_w + gap), top, card_w, height, theme)


def _layout_card_and_code(slide, card_sections: list[dict], code_sections: list[dict],
                          top: float, height: float, theme: Theme):
    ml = theme.margin_left
    gap = theme.card_gap
    total_w = theme.content_width

    if card_sections and code_sections:
        card_w = total_w * 0.45
        _render_card(slide, card_sections[0], ml, top, card_w, height, theme)
        _add_code_block(slide, ml + card_w + gap, top, total_w - card_w - gap, height,
                        code_sections[0]["content"][0]["text"], theme)
        return

    if code_sections:
        _add_code_block(slide, ml, top, total_w, height,
                        code_sections[0]["content"][0]["text"], theme)
        return

    if not card_sections:
        return

    sec = card_sections[0]
    code_items = [c for c in sec["content"] if c["type"] == "code"]
    non_code = [c for c in sec["content"] if c["type"] != "code"]

    if not code_items:
        return

    if non_code:
        card_w = total_w * 0.45
        _render_card(slide, {**sec, "content": non_code}, ml, top, card_w, height, theme)
        _add_code_block(slide, ml + card_w + gap, top, total_w - card_w - gap, height,
                        code_items[0]["text"], theme)
    else:
        _add_code_block(slide, ml, top, total_w, height, code_items[0]["text"], theme)


def _render_card(slide, section: dict, left: float, top: float,
                 width: float, height: float, theme: Theme):
    accent = section.get("accent")
    dark_card = _is_dark_color(theme.card_bg)
    _add_card(slide, left, top, width, height, theme, accent, dark_card=dark_card)

    pad = min(theme.card_padding, width * 0.08)
    inner_left = left + pad
    inner_width = width - pad * 2
    y = top + pad

    if section.get("heading"):
        if dark_card:
            heading_color = theme.text_on_dark
        else:
            heading_color = theme.accent(accent) if accent else theme.text_heading
        card_title_size = min(theme.font_size("card_title", 16), 16)
        th = _estimate_text_height(section["heading"], card_title_size, inner_width, bold=True)
        _add_textbox(slide, inner_left, y, inner_width, th,
                     section["heading"], card_title_size,
                     heading_color, theme, bold=True)
        y += th + 0.1

    text_color = theme.text_on_dark if dark_card else theme.text_body
    heading_text_color = theme.text_on_dark if dark_card else theme.text_heading
    body_size = min(theme.font_size("body", 13), 14)
    small_size = min(theme.font_size("small", 12), 13)

    for item in section.get("content", []):
        if item["type"] in ("text", "bold"):
            is_bold = item["type"] == "bold"
            size = body_size if is_bold else small_size
            color = heading_text_color if is_bold else text_color
            th = _estimate_text_height(item["text"], size, inner_width, bold=is_bold)
            _add_textbox(slide, inner_left, y, inner_width, th,
                         item["text"], size, color, theme, bold=is_bold)
            y += th + 0.04

    bullets = [c for c in section.get("content", []) if c["type"] == "bullet"]
    if bullets:
        remaining_h = max(0.3, height - (y - top) - pad)
        bullet_size = min(body_size, 13)
        total_bullet_lines = sum(
            max(1, -(-len(b["text"]) // max(1, int(inner_width * 72 / (bullet_size * 0.55)))))
            for b in bullets
        )
        estimated_bullet_h = total_bullet_lines * bullet_size * 1.4 / 72
        if estimated_bullet_h > remaining_h and remaining_h > 0.3:
            bullet_size = max(9, int(bullet_size * remaining_h / estimated_bullet_h))
        _add_bullet_list(slide, inner_left, y, inner_width, remaining_h,
                         bullets, theme,
                         font_size=bullet_size,
                         text_color=text_color,
                         bullet_color=theme.accent(accent) if accent else None)



def _render_takeaway(slide, text: str, theme: Theme):
    ml = theme.margin_left
    comp = theme.comp.get("callout", {})
    h = comp.get("height", 0.6)
    top = theme.slide_height - h - comp.get("bottom_margin", 0.4)

    dark_card = _is_dark_color(theme.card_bg)
    _add_card(slide, ml, top, theme.content_width, h, theme,
              comp.get("accent", "dark"), dark_card=dark_card)

    pad = theme.card_padding
    text_color = theme.text_on_dark if dark_card else theme.text_heading
    _add_textbox(slide, ml + pad, top + 0.1, theme.content_width - pad * 2, h - 0.2,
                 text, comp.get("font_size", 12),
                 text_color, theme, bold=comp.get("bold", True))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_presentation(md_path: str, theme_path: str | None, output_path: str):
    md_text = Path(md_path).read_text()
    md_dir = Path(md_path).parent

    frontmatter, body = parse_frontmatter(md_text)

    if theme_path is None:
        theme_ref = frontmatter.get("theme")
        theme_path = str(md_dir / theme_ref) if theme_ref else None

    if theme_path is None:
        default = Path(__file__).parent.parent / "references" / "default-theme.json"
        if default.exists():
            theme_path = str(default)
        else:
            print("Error: No theme specified and no default-theme.json found", file=sys.stderr)
            sys.exit(1)

    theme = Theme.load(theme_path)

    prs = Presentation()
    prs.slide_width = Inches(theme.slide_width)
    prs.slide_height = Inches(theme.slide_height)

    for slide_data in parse_slides(body):
        mods = slide_data.get("mods", [])
        if "title" in mods:
            _build_title_slide(prs, slide_data, theme, frontmatter, str(md_dir))
        elif "dark" in mods:
            _build_dark_slide(prs, slide_data, theme)
        else:
            _build_content_slide(prs, slide_data, theme)

    prs.save(output_path)
    print(f"Saved {len(prs.slides)} slides to {output_path}")


def main():
    parser = argparse.ArgumentParser(description="Build a PPTX from structured markdown")
    parser.add_argument("markdown", help="Path to the structured markdown file")
    parser.add_argument("--theme", "-t", help="Path to theme.json")
    parser.add_argument("--output", "-o", default="presentation.pptx", help="Output PPTX path")
    args = parser.parse_args()
    build_presentation(args.markdown, args.theme, args.output)


if __name__ == "__main__":
    main()
